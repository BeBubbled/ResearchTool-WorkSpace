#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Export a PowerPoint (.pptx) with an R x C grid of images (no extra borders).
- Cross-platform (Windows/macOS/Linux)
- Fixed per-cell size (cm), fixed gaps/margins (cm/pt/px)
- 'fit' (letterbox) or 'fill' (crop) placement inside each cell
- Slide size auto-fits the grid unless --origin-cm is given (then slide size unchanged)

pip install python-pptx Pillow

Example:
python sort_images_ppt.py --images "Images" --rows 3 --cols 5 --cell-size-cm 5 --gap-px 4 --margin-cm 1 --cell-aspect 1:1 --fit fit --sort numeric --out grid_4x4_5cm.pptx
"""

import argparse
import glob
import os
import re
import sys
from typing import List

from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt

# --- Unit helpers ---
PT_PER_INCH = 72.0
CM_PER_INCH = 2.54

def cm_to_pt(cm: float) -> float:
    return cm / CM_PER_INCH * PT_PER_INCH

def px_to_pt(px: float, dpi: float = 96.0) -> float:
    return px / dpi * PT_PER_INCH

def pt_to_cm(pt: float) -> float:
    return pt / PT_PER_INCH * CM_PER_INCH

def parse_aspect(s: str) -> float:
    parts = s.split(":")
    if len(parts) != 2:
        raise argparse.ArgumentTypeError("Aspect must be like '1:1' or '4:3'.")
    w, h = float(parts[0]), float(parts[1])
    if h == 0:
        raise argparse.ArgumentTypeError("Aspect denominator cannot be 0.")
    return w / h

# --- CLI ---
def parse_args():
    p = argparse.ArgumentParser(
        description="Export a PPTX with images arranged in a fixed-size R x C grid (no borders)."
    )
    p.add_argument("--images", nargs="+", required=True,
                   help="Image paths or glob(s), e.g., path\\to\\*.png")
    p.add_argument("--rows", type=int, required=True)
    p.add_argument("--cols", type=int, required=True)

    g = p.add_mutually_exclusive_group(required=True)
    g.add_argument("--cell-size-cm", type=float,
                   help="Square cell (width==height) in cm, e.g., 5 for 5x5 cm.")
    p.add_argument("--cell-width-cm", type=float,
                   help="Cell width in cm (if not using --cell-size-cm).")
    p.add_argument("--cell-height-cm", type=float,
                   help="Cell height in cm (if not using --cell-size-cm).")

    p.add_argument("--cell-aspect", type=parse_aspect,
                   help="Expected display aspect W:H (e.g., 1:1). Only checks; no forcing.")

    gap = p.add_mutually_exclusive_group()
    gap.add_argument("--gap-cm", type=float, help="Gap between cells in cm.")
    gap.add_argument("--gap-pt", type=float, help="Gap between cells in points.")
    gap.add_argument("--gap-px", type=float, help="Gap between cells in pixels (96 DPI assumed).")

    margin = p.add_mutually_exclusive_group()
    margin.add_argument("--margin-cm", type=float, help="Outer margin in cm.")
    margin.add_argument("--margin-pt", type=float, help="Outer margin in points.")
    margin.add_argument("--margin-px", type=float, help="Outer margin in pixels (96 DPI assumed).")

    p.add_argument("--fit", choices=["fit", "fill"], default="fit",
                   help="'fit' keeps entire image (letterbox). 'fill' crops to fill.")
    p.add_argument("--sort", choices=["name", "mtime", "none","numeric"], default="name")
    p.add_argument("--reverse", action="store_true")

    p.add_argument("--origin-cm", nargs=2, type=float, metavar=("XCM", "YCM"),
                   help="Grid top-left origin (cm) on the slide. If omitted, slide auto-sizes to grid.")

    p.add_argument("--out", type=str, default="grid_output.pptx",
                   help="Output .pptx path. Default: grid_output.pptx")
    return p.parse_args()

# --- Helpers ---
def expand_images(args_images: List[str]) -> List[str]:
    """
    支持三种输入:
    1) 带通配符的模式: "./Images/*.png"
    2) 目录: "Images"
    3) 单个文件: "xxx.png"
    """
    paths: List[str] = []

    for term in args_images:
        # 1. 带通配符: 用 glob 展开
        if any(ch in term for ch in ["*", "?", "["]):
            paths.extend(glob.glob(term))
            continue

        # 2. 目录: 递归找常见图片后缀
        if os.path.isdir(term):
            for root, _, files in os.walk(term):
                for f in files:
                    if f.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".webp")):
                        paths.append(os.path.join(root, f))
            continue

        # 3. 文件: 直接加入
        if os.path.isfile(term):
            paths.append(term)
            continue

        # 4. 既不是文件也不是目录
        print(f"WARNING: '{term}' is not a file or directory, ignored.", file=sys.stderr)

    return paths


def sort_images(imgs: List[str], mode: str) -> List[str]:
    """
    mode:
      - none: 不排序
      - name: 字典序
      - mtime: 按修改时间
      - numeric: 按文件名中的数字排序 (最后一段数字)
    """
    if mode == "none":
        return imgs

    if mode == "name":
        return sorted(
            imgs,
            key=lambda p: (os.path.dirname(p), os.path.basename(p).lower())
        )

    if mode == "mtime":
        return sorted(imgs, key=lambda p: os.path.getmtime(p))

    if mode == "numeric":
        def extract_number(fname: str) -> int:
            # 取 basename 中最后一段连续数字:
            # 例如: "frame_23_1_10.png" -> 10
            base = os.path.basename(fname)
            m = re.search(r"(\d+)(?!.*\d)", base)
            if m:
                return int(m.group(1))
            # 找不到数字就排在最前或最后，这里给一个特别小的值
            return -1

        return sorted(imgs, key=extract_number)

    return imgs


def compute_units(args):
    # Cell size (cm)
    if args.cell_size_cm is not None:
        cw_cm = ch_cm = float(args.cell_size_cm)
    else:
        if args.cell_width_cm is None or args.cell_height_cm is None:
            raise SystemExit("Provide both --cell-width-cm and --cell-height-cm, or use --cell-size-cm.")
        cw_cm = float(args.cell_width_cm)
        ch_cm = float(args.cell_height_cm)

    # Aspect check (optional)
    if args.cell_aspect is not None:
        ratio = cw_cm / max(ch_cm, 1e-6)
        if abs(ratio - args.cell_aspect) > 1e-3:
            print(f"WARNING: Cell aspect {ratio:.3f} != expected {args.cell_aspect:.3f}.", file=sys.stderr)

    # Gap
    if args.gap_cm is not None:
        gap_cm = float(args.gap_cm)
    elif args.gap_pt is not None:
        gap_cm = pt_to_cm(float(args.gap_pt))
    elif args.gap_px is not None:
        gap_cm = pt_to_cm(px_to_pt(float(args.gap_px)))
    else:
        gap_cm = 0.2  # default 0.2 cm

    # Margin
    if args.margin_cm is not None:
        margin_cm = float(args.margin_cm)
    elif args.margin_pt is not None:
        margin_cm = pt_to_cm(float(args.margin_pt))
    elif args.margin_px is not None:
        margin_cm = pt_to_cm(px_to_pt(float(args.margin_px)))
    else:
        margin_cm = 1.0  # default 1.0 cm

    return cw_cm, ch_cm, gap_cm, margin_cm

def grid_size_cm(rows: int, cols: int, cw_cm: float, ch_cm: float, gap_cm: float, margin_cm: float):
    grid_w_cm = cols * cw_cm + (cols - 1) * gap_cm + 2 * margin_cm
    grid_h_cm = rows * ch_cm + (rows - 1) * gap_cm + 2 * margin_cm
    return grid_w_cm, grid_h_cm

def add_image_in_cell_cm(slide, img_path: str,
                         cell_left_cm: float, cell_top_cm: float,
                         cell_w_cm: float, cell_h_cm: float,
                         fit_mode: str = "fit"):
    """Place an image inside a cell rectangle (all sizes in cm)."""
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    img_ar = w_px / max(h_px, 1)
    cell_ar = cell_w_cm / max(cell_h_cm, 1e-9)

    if fit_mode == "fit":
        # keep entire image
        if img_ar >= cell_ar:
            pic_w_cm = cell_w_cm
            pic_h_cm = pic_w_cm / img_ar
        else:
            pic_h_cm = cell_h_cm
            pic_w_cm = pic_h_cm * img_ar
    else:  # fill (crop)
        if img_ar >= cell_ar:
            # image wider; match height
            pic_h_cm = cell_h_cm
            pic_w_cm = pic_h_cm * img_ar
        else:
            # image taller; match width
            pic_w_cm = cell_w_cm
            pic_h_cm = pic_w_cm / img_ar

    left_cm = cell_left_cm + (cell_w_cm - pic_w_cm) / 2.0
    top_cm  = cell_top_cm  + (cell_h_cm - pic_h_cm) / 2.0

    slide.shapes.add_picture(
        img_path,
        left=Cm(left_cm),
        top=Cm(top_cm),
        width=Cm(pic_w_cm),
        height=Cm(pic_h_cm),
    )

def main():
    args = parse_args()
    imgs = expand_images(args.images)
    if not imgs:
        print("No images found.", file=sys.stderr)
        sys.exit(2)

    imgs = sort_images(imgs, args.sort)
    if args.reverse:
        imgs = list(reversed(imgs))
    print("[INFO] Sorted images (after sort / reverse):")
    for i, p in enumerate(imgs):
        print(f"  {i:03d}: {os.path.basename(p)}")
    print(f"[INFO] Total images used (before trim/pad): {len(imgs)}")

    needed = args.rows * args.cols
    if len(imgs) < needed:
        print(f"WARNING: Need {needed} images, but only {len(imgs)} provided. Filling remaining cells with last image.", file=sys.stderr)
        if imgs:
            imgs = imgs + [imgs[-1]] * (needed - len(imgs))
        else:
            print("No images at all.", file=sys.stderr)
            sys.exit(2)
    else:
        imgs = imgs[:needed]

    cw_cm, ch_cm, gap_cm, margin_cm = compute_units(args)
    grid_w_cm, grid_h_cm = grid_size_cm(args.rows, args.cols, cw_cm, ch_cm, gap_cm, margin_cm)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank

    if args.origin_cm:
        # Keep default slide size; place grid at the given origin
        slide = prs.slides.add_slide(blank_layout)
        origin_left_cm = float(args.origin_cm[0])
        origin_top_cm  = float(args.origin_cm[1])
    else:
        # Auto-size slide to the exact grid size
        prs.slide_width  = Cm(grid_w_cm)
        prs.slide_height = Cm(grid_h_cm)
        slide = prs.slides.add_slide(blank_layout)
        origin_left_cm = 0.0
        origin_top_cm  = 0.0

    idx = 0
    for r in range(args.rows):
        for c in range(args.cols):
            left_cm = origin_left_cm + margin_cm + c * (cw_cm + gap_cm)
            top_cm  = origin_top_cm  + margin_cm + r * (ch_cm + gap_cm)
            add_image_in_cell_cm(
                slide, imgs[idx],
                cell_left_cm=left_cm, cell_top_cm=top_cm,
                cell_w_cm=cw_cm, cell_h_cm=ch_cm,
                fit_mode=args.fit
            )
            idx += 1

    prs.save(args.out)
    print(f"Saved: {args.out}")
    print("Open the PPTX, select shapes (Ctrl/Cmd+A), Group (optional), then copy to your target deck.")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print("ERROR:", e, file=sys.stderr)
        sys.exit(1)