#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
stack_cli.py

- images：线性横/纵排列图片 → PPTX（每图独立形状，默认间隔5px、默认1pt黑边）
- videos：rows×cols 网格拼接视频（白底黑字，默认间隔5px、外侧留白5px）
  readme.txt 严格格式：
    前 N 行（N=rows*cols）：每格标题
    接着两行空行
    若 direction=h：再给 rows 行 → 行 caption（可用空行表示该行空标题）
    若 direction=v：再给 cols 行 → 列 caption（可用空行表示该列空标题）

依赖：
  pip install python-pptx Pillow
  并安装 ffmpeg/ffprobe 到 PATH

注意：
  - 未提供 --fontfile 时，drawtext 使用 font='Arial'
  - 若输出文件已存在，自动追加 _1, _2, ...；不覆盖
"""

import argparse
import glob
import sys
import subprocess
from pathlib import Path
from typing import List, Tuple, Optional
import tempfile
import os


PT_PER_INCH = 72.0
CM_PER_INCH = 2.54
DEFAULT_DPI = 96.0

def px_to_cm(px: float, dpi: float = DEFAULT_DPI) -> float:
    return (px / dpi) * 2.54

def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem, suf, parent = path.stem, path.suffix, path.parent
    i = 1
    while True:
        p = parent / f"{stem}_{i}{suf}"
        if not p.exists():
            return p
        i += 1

def natural_sort_key(p: str):
    q = Path(p)
    return (q.name.lower(), str(q).lower())

def _escape_path_for_drawtext(p: str) -> str:
    # ffmpeg 的 drawtext 选项用冒号分隔，路径里的冒号/反斜杠/引号需要转义
    return p.replace("\\", r"\\").replace(":", r"\:").replace("'", r"\'")

# ----------------- images 子命令（PPTX） -----------------
from PIL import Image
from PIL import ImageFont, ImageDraw
from pptx import Presentation
from pptx.util import Cm, Pt

def cmd_images(args: argparse.Namespace) -> int:
    files: List[str] = []
    for term in args.files:
        if any(ch in term for ch in ["*", "?", "["]):
            files.extend(glob.glob(term))
        else:
            files.append(term)
    files = [f for f in files if Path(f).is_file()]
    if not files:
        print("No image files found.", file=sys.stderr)
        return 2

    files.sort(key=natural_sort_key)

    sizes_px = []
    for f in files:
        try:
            with Image.open(f) as im:
                sizes_px.append((im.width, im.height))
        except Exception as e:
            print(f"Failed to open image: {f} ({e})", file=sys.stderr)
            return 2

    gap_cm = px_to_cm(args.gap_px)
    border_pt = max(0.0, float(args.border_pt))

    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    if args.direction == "horizontal":
        base_h = sizes_px[0][1]
        widths_cm, heights_cm = [], []
        for (w, h) in sizes_px:
            scale = base_h / h
            widths_cm.append(px_to_cm(w * scale))
            heights_cm.append(px_to_cm(base_h))
        total_w = sum(widths_cm) + gap_cm * (len(files) - 1)
        total_h = max(heights_cm)
        prs.slide_width = Cm(total_w)
        prs.slide_height = Cm(total_h)

        x = 0.0
        y = 0.0
        for i, f in enumerate(files):
            wcm, hcm = widths_cm[i], heights_cm[i]
            pic = slide.shapes.add_picture(f, left=Cm(x), top=Cm(y), width=Cm(wcm), height=Cm(hcm))
            if border_pt > 0:
                pic.line.width = Pt(border_pt)
            x += wcm + gap_cm
    else:
        base_w = sizes_px[0][0]
        widths_cm, heights_cm = [], []
        for (w, h) in sizes_px:
            scale = base_w / w
            widths_cm.append(px_to_cm(base_w))
            heights_cm.append(px_to_cm(h * scale))
        total_w = max(widths_cm)
        total_h = sum(heights_cm) + gap_cm * (len(files) - 1)
        prs.slide_width = Cm(total_w)
        prs.slide_height = Cm(total_h)

        x = 0.0
        y = 0.0
        for i, f in enumerate(files):
            wcm, hcm = widths_cm[i], heights_cm[i]
            pic = slide.shapes.add_picture(f, left=Cm(x), top=Cm(y), width=Cm(wcm), height=Cm(hcm))
            if border_pt > 0:
                pic.line.width = Pt(border_pt)
            y += hcm + gap_cm

    out_path = unique_path(Path(args.out))
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    return 0

# ----------------- videos 子命令（ffmpeg） -----------------

VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".avi"}

def is_video(p: Path) -> bool:
    return p.suffix.lower() in VIDEO_EXTS

def probe_video_size(p: Path) -> Optional[Tuple[int, int]]:
    cmd = [
        "ffprobe", "-v", "error", "-select_streams", "v:0",
        "-show_entries", "stream=width,height", "-of", "csv=p=0:s=x", str(p)
    ]
    try:
        out = subprocess.check_output(cmd, stderr=subprocess.DEVNULL).decode().strip()
        if out:
            w, h = out.split("x")
            return int(w), int(h)
    except Exception:
        return None
    return None

def probe_video_fps(p: Path) -> Optional[float]:
    cmd = [
        "ffprobe", "-v", "error", "-select_streams", "v:0",
        "-show_entries", "stream=r_frame_rate", "-of", "csv=p=0", str(p)
    ]
    try:
        out = subprocess.check_output(cmd, stderr=subprocess.DEVNULL).decode().strip()
        if out and "/" in out:
            num, den = out.split("/")
            den = float(den) if float(den) != 0 else 1.0
            return float(num) / den
        elif out:
            return float(out)
    except Exception:
        return None
    return None

def escape_drawtext(s: str) -> str:
    return s.replace("\\", r"\\").replace(":", r"\:").replace("'", r"\'").replace("%", r"\%")

def _load_pil_font(fontfile: str, fontsize: int) -> ImageFont.FreeTypeFont:
    try:
        if fontfile:
            return ImageFont.truetype(fontfile, fontsize)
    except Exception:
        pass
    # Fallbacks
    try:
        return ImageFont.truetype("Arial.ttf", fontsize)
    except Exception:
        return ImageFont.load_default()

def wrap_text_to_width(text: str, fontfile: str, fontsize: int, max_width_px: int) -> List[str]:
    if max_width_px <= 0 or not text:
        return [text]
    font = _load_pil_font(fontfile, fontsize)
    img = Image.new("RGB", (max_width_px, fontsize*4), color=(255,255,255))
    draw = ImageDraw.Draw(img)

    # Decide tokenization: prefer word-based; if no spaces (e.g., CJK), fall back to char-based
    if any(ch.isspace() for ch in text):
        tokens = text.split()
        sep = " "
    else:
        tokens = list(text)
        sep = ""

    lines: List[str] = []
    current: List[str] = []
    for tok in tokens:
        candidate = (sep.join(current + [tok])).strip()
        w, _ = draw.textsize(candidate, font=font)
        if w <= max_width_px or not current:
            current.append(tok)
        else:
            lines.append(sep.join(current).strip())
            current = [tok]
    if current:
        lines.append(sep.join(current).strip())
    if not lines:
        lines = [text]
    return lines

def build_multiline_text(text: str, fontfile: str, fontsize: int, max_width_px: int) -> str:
    lines = wrap_text_to_width(text, fontfile, fontsize, max_width_px)
    escaped_lines = [escape_drawtext(ln) for ln in lines]
    joined = "{NL}".join(escaped_lines)
    return joined.replace("{NL}", r"\n")

def measure_multiline_height(text: str, fontfile: str, fontsize: int, max_width_px: int) -> Tuple[int, str]:
    """Return (pixel_height, wrapped_text_with_newlines)."""
    lines = wrap_text_to_width(text, fontfile, fontsize, max_width_px)
    joined = "\n".join(lines)
    font = _load_pil_font(fontfile, fontsize)
    # Create a small canvas; bbox does not depend on canvas actual size
    img = Image.new("RGB", (max_width_px, fontsize * (len(lines) + 2)), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    try:
        # Pillow >= 8: multiline_textbbox available
        bbox = draw.multiline_textbbox((0, 0), joined, font=font, spacing=0, align="left")
    except Exception:
        bbox = draw.textbbox((0, 0), joined, font=font)
    height = max(0, bbox[3] - bbox[1])
    return height, joined

def reorder_for_direction(items: List[str], rows: int, cols: int, to_col_major: bool) -> List[str]:
    if not to_col_major:
        return items[:]
    grid = [items[r*cols:(r+1)*cols] for r in range(rows)]
    out = []
    for c in range(cols):
        for r in range(rows):
            out.append(grid[r][c])
    return out

def parse_readme_strict(readme_path: Path, n_items: int, rows: int, cols: int, direction: str):
    """严格解析：
       前 n_items 行为每格标题；随后两行空行；接着
       - direction='h' → 读取 rows 行为行 caption
       - direction='v' → 读取 cols 行为列 caption
       若 readme 不存在或行数不足：用空串补齐；多余行忽略。
    """
    per_item = ["" for _ in range(n_items)]
    row_caps = ["" for _ in range(rows)]
    col_caps = ["" for _ in range(cols)]

    if not readme_path.is_file():
        return per_item, row_caps, col_caps

    with open(readme_path, "r", encoding="utf-8") as f:
        lines = [ln.rstrip("\n") for ln in f.readlines()]

    # 1) 每格标题
    for i in range(n_items):
        if i < len(lines):
            per_item[i] = lines[i]

    # 2) 两个空行分隔
    idx = n_items
    # 宽容处理：允许中间存在额外空行或缺一行，但我们严格“至少两行空行”
    # 不强制报错，只是按规则尽量前进
    empties = 0
    while idx < len(lines) and empties < 2:
        if lines[idx].strip() == "":
            empties += 1
        else:
            empties = 0
        idx += 1

    # 3) 根据方向读取行/列 caption
    if direction == "h":
        for r in range(rows):
            j = idx + r
            if j < len(lines):
                row_caps[r] = lines[j]
    else:  # 'v'
        for c in range(cols):
            j = idx + c
            if j < len(lines):
                col_caps[c] = lines[j]

    return per_item, row_caps, col_caps

def check_aspect_consistency(sizes: List[Tuple[int,int]], tol: float = 1e-3) -> bool:
    w0, h0 = sizes[0]
    r0 = w0 / h0
    for (w, h) in sizes[1:]:
        if h == 0:
            return False
        if abs((w / h) - r0) > tol:
            return False
    return True

def build_xstack_layout_with_bands(
    cols: int, rows: int,
    tile_w: int, tile_h: int,
    gap: int, outer: int,
    title_band_px: int, rowcap_band_px: int
):
    """返回 xstack layout、tile 左上 (x,y)、行区域 (x, row_top_y, row_w)，不含右/下 pad。"""
    row_span = title_band_px + tile_h + rowcap_band_px
    total_w = cols * tile_w + (cols - 1) * gap + 2 * outer
    total_h = rows * row_span + (rows - 1) * gap + 2 * outer

    row_w = cols * tile_w + (cols - 1) * gap

    tile_xy = []
    for r in range(rows):
        row_y_top = outer + r * (row_span + gap)
        for c in range(cols):
            x = outer + c * (tile_w + gap)
            y = row_y_top + title_band_px
            tile_xy.append((x, y))

    row_xyw = []
    for r in range(rows):
        row_y_top = outer + r * (row_span + gap)
        row_xyw.append((outer, row_y_top, row_w))

    layout = "|".join(f"{x}_{y}" for (x, y) in tile_xy)
    return layout, total_w, total_h, tile_xy, row_xyw

def cmd_videos(args: argparse.Namespace) -> int:
    folder = Path(args.dir)
    if not folder.is_dir():
        print("Invalid --dir.", file=sys.stderr)
        return 2

    direction = args.mode.lower()  # 'h' or 'v'
    readme = folder / "readme.txt"  # 默认

    vids = sorted([str(p) for p in folder.iterdir() if p.is_file() and is_video(p)], key=natural_sort_key)
    if not vids:
        print("No videos found in --dir.", file=sys.stderr)
        return 2

    need = args.rows * args.cols
    vids = vids[:need] if len(vids) >= need else vids + [vids[-1]] * (need - len(vids))

    # 严格解析 readme
    per_item_titles_raw, row_caps, col_caps = parse_readme_strict(
        readme, n_items=need, rows=args.rows, cols=args.cols, direction=direction
    )

    # 探测尺寸并取最高分辨率
    sizes = []
    for v in vids:
        s = probe_video_size(Path(v))
        if s is None:
            print(f"ffprobe failed on {v}", file=sys.stderr)
            return 2
        sizes.append(s)

    if not check_aspect_consistency(sizes):
        print("Warning: video aspect ratios differ; forcing scale to max resolution.")

    areas = [w*h for (w, h) in sizes]
    base_idx = max(range(len(areas)), key=lambda i: areas[i])
    base_w, base_h = sizes[base_idx]
    fps = probe_video_fps(Path(vids[0])) or 30.0

    # 方向：重排视频与“每格标题”到行主序 tile 顺序
    to_col_major = (direction == "v")
    vids = reorder_for_direction(vids, args.rows, args.cols, to_col_major)
    per_item_titles = reorder_for_direction(per_item_titles_raw, args.rows, args.cols, to_col_major)

    # 默认参数
    gap        = int(getattr(args, "gap_px", 5))
    outer      = int(getattr(args, "outer_border_px", 5))
    title_band = max(0, int(getattr(args, "title_band_px", 40)))
    # 行/列 band：按方向启用默认值；另一种保持 0（不占带）
    row_band_default = 48 if direction == "h" else 0
    col_band_default = 48 if direction == "v" else 0
    row_band   = max(0, int(getattr(args, "rowcap_band_px", row_band_default)))
    col_band   = max(0, int(getattr(args, "colcap_band_px", col_band_default)))
    title_fs   = int(getattr(args, "title_fontsize", 22))
    rowcap_fs  = int(getattr(args, "rowcap_fontsize", 26))
    colcap_fs  = int(getattr(args, "colcap_fontsize", 26))
    keep_audio = getattr(args, "keep_audio", "first")
    fontfile   = getattr(args, "fontfile", "")


    # === 在这里插入你的“增强测高”代码块 ===
    # （增强）h 模式下：按“受限宽度”测量行 caption 高度，必要时增大 row_band
    if direction == "h" and row_band > 0:
        # 行总宽（含 gap）：row_w_full = cols*base_w + (cols-1)*gap
        row_w_full = args.cols * base_w + (args.cols - 1) * gap

        # 受限区域：max(2, cols-3) 个 tile 宽，不含 gap（按你的要求）
        area_tiles = max(2, args.cols - 3)
        area_w = area_tiles * base_w
        # 留左右各 5px 内边距
        max_rowcap_w = max(1, area_w - 10)

        max_h = 0
        for r in range(args.rows):
            cap_text = row_caps[r] if r < len(row_caps) else ""
            h_px, _ = measure_multiline_height(cap_text, fontfile, rowcap_fs, max_rowcap_w)
            if h_px > max_h:
                max_h = h_px
        row_band = max(row_band, max_h)

    # 1) scale
    filter_parts = []
    in_labels = []
    for i, v in enumerate(vids):
        chain = f"[{i}:v]scale={base_w}:{base_h}:flags=lanczos[v{i}]"
        filter_parts.append(chain)
        in_labels.append(f"[v{i}]")

    # 2) 行带（row_band）只在 h 模式有效；v 模式 row_band=0
    layout, total_w, total_h, tile_xy, row_xyw = build_xstack_layout_with_bands(
        cols=args.cols, rows=args.rows,
        tile_w=base_w, tile_h=base_h,
        gap=gap, outer=outer,
        title_band_px=title_band, rowcap_band_px=row_band
    )

    # 3) xstack（白底）
    inputs = "".join(in_labels)
    filter_parts.append(
        f"{inputs}xstack=inputs={len(in_labels)}:layout={layout}:fill=white[grid]"
    )

    # 3.5) pad 右/下并补偶数；高度需考虑列带（col_band）及底部 outer
    # 当前 ih = 上外边距 + rows*(title + tile + row_band) + (rows-1)*gap + 下外边距
    # 若 direction='v'，再加列带；列带高度需容纳换行后的文字
    pad_w_val = total_w + outer  # 右侧追加 outer 像素留白

    if direction == "v" and col_band > 0:
        # 预先测量每列 caption 在最大宽度 base_w-10 下的像素高度，取最大值
        max_colcap_w = max(1, base_w - 10)
        max_cap_h = 0
        for c in range(args.cols):
            cap_text = col_caps[c] if c < len(col_caps) else ""
            h_px, _ = measure_multiline_height(cap_text, fontfile, colcap_fs, max_colcap_w)
            if h_px > max_cap_h:
                max_cap_h = h_px
        # 至少保留原有 col_band，若文本更高则按文本高度
        effective_col_band = max(col_band, max_cap_h)
        pad_h_val = total_h + effective_col_band
        computed_col_band = effective_col_band
    else:
        pad_h_val = total_h
        computed_col_band = 0
    filter_parts.append(
        f"[grid]pad=w={pad_w_val}:h={pad_h_val}:x=0:y=0:color=white[grid]"
    )
    # 最后保证偶数尺寸，满足编码器需求
    filter_parts.append("[grid]scale=ceil(iw/2)*2:ceil(ih/2)*2[grid]")

    # 4) 每格标题（行上方带内，黑字、半透明白底）
    for i, (x, y) in enumerate(tile_xy):
        if title_band <= 0:
            continue
        # 标题宽度：与 tile 等宽，左右各留 5px
        max_title_w = max(1, base_w - 10)
        title = build_multiline_text(per_item_titles[i], fontfile, title_fs, max_title_w)
        if not title.strip():
            continue
        draw = "[grid]drawtext="
        if fontfile:
            draw += f"fontfile='{fontfile}':"
        else:
            draw += "font='Arial':"
        draw += (
            f"text='{title}':"
            f"x={x}+5+( {base_w} - 10 - text_w )/2:"
            f"y={y - title_band}+( {title_band} - text_h )/2:"
            f"fontsize={title_fs}:fontcolor=black:"
            f"box=1:boxcolor=white@0.5:boxborderw=4"
            f"[grid]"
        )
        filter_parts.append(draw)

    # 5a) 行 caption（h 模式）：行下方带内，使用“受限宽度”并居中；用 textfile 确保换行
    if direction == "h" and row_band > 0:
        for r in range(args.rows):
            row_x, row_y_top, row_w = row_xyw[r]

            # 受限区域宽度：max(2, cols-3) * base_w（不含 gap）
            area_tiles = max(2, args.cols - 3)
            area_w = area_tiles * base_w

            # 将该受限区域水平居中到整行（整行宽度 row_w 含 gap）
            area_left = row_x + (row_w - area_w) / 2.0

            # 文本最大宽度 = 受限区域宽度 - 10（左右各留 5px）
            max_rowcap_w = max(1, int(area_w - 10))

            raw_text = row_caps[r] if r < len(row_caps) else ""
            # 先在受限宽度下断行
            _, wrapped_raw = measure_multiline_height(raw_text, fontfile, rowcap_fs, max_rowcap_w)

            # 写入临时文件（真实换行符）
            tf = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
            tf.write(wrapped_raw.encode("utf-8"))
            tf.close()
            path_escaped = _escape_path_for_drawtext(tf.name)

            tile_top_y = row_y_top + title_band
            tile_bottom_y = tile_top_y + base_h

            draw = "[grid]drawtext="
            if fontfile:
                draw += f"fontfile='{fontfile}':"
            else:
                draw += "font='Arial':"
            draw += (
                f"textfile='{path_escaped}':"
                f"x={area_left}+5+( {max_rowcap_w} - text_w )/2:"
                f"y={tile_bottom_y}+( {row_band} - text_h )/2:"
                f"fontsize={rowcap_fs}:fontcolor=black:"
                f"line_spacing=0:"
                f"box=1:boxcolor=white@0.5:boxborderw=6"
                f"[grid]"
            )
            filter_parts.append(draw)

    # 5b) 列 caption（v 模式）：整幅底部新增 col_band 内，按列居中
    if direction == "v" and col_band > 0:
        col_band_top = total_h  # 列带紧贴网格底部
        max_colcap_w = max(1, base_w - 10)

        # 为每列写入一个临时文本文件（内容已按宽度换行）
        temp_files = []
        try:
            for c in range(args.cols):
                raw_text = col_caps[c] if c < len(col_caps) else ""
                # 使用与测量一致的换行策略
                _, wrapped_raw = measure_multiline_height(raw_text, fontfile, colcap_fs, max_colcap_w)

                # 写入临时文件（真实换行符）
                tf = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
                tf.write(wrapped_raw.encode("utf-8"))
                tf.close()
                temp_files.append(tf.name)

                col_left_x = outer + c * (base_w + gap)
                path_escaped = _escape_path_for_drawtext(tf.name)

                draw = "[grid]drawtext="
                if fontfile:
                    draw += f"fontfile='{fontfile}':"
                else:
                    draw += "font='Arial':"
                draw += (
                    f"textfile='{path_escaped}':"
                    f"x={col_left_x}+5+( {base_w} - 10 - text_w )/2:"
                    f"y={col_band_top}+( {computed_col_band} - text_h )/2:"
                    f"fontsize={colcap_fs}:fontcolor=black:"
                    f"line_spacing=0:"
                    f"box=1:boxcolor=white@0.5:boxborderw=6"
                    f"[grid]"
                )
                filter_parts.append(draw)
        finally:
            # 提示：是否删除临时文件由你决定；多数情况下可以保留到进程结束后由系统清理
            # 如需主动删除，可在 ffmpeg 运行完后 os.unlink
            pass

    filter_complex = ";".join(filter_parts)

    # 6) ffmpeg 命令
    cmd = ["ffmpeg", "-y"]
    for v in vids:
        cmd += ["-i", v]
    cmd += ["-filter_complex", filter_complex, "-map", "[grid]"]
    if keep_audio == "first":
        cmd += ["-map", "0:a?", "-c:a", "aac", " -b:a", "192k"]  # 空格注意
        # 修正参数列表的空格问题：
        cmd[-2:] = ["-b:a", "192k"]
    else:
        cmd += ["-an"]
    cmd += [
        "-r", f"{fps:.3f}",
        "-c:v", "libx264", "-pix_fmt", "yuv420p",
        "-crf", "18", "-preset", "medium",
        "-movflags", "+faststart",
    ]
    out_path = unique_path(Path(args.out))
    cmd += [str(out_path)]

    print("Running ffmpeg:")
    shown = " ".join(cmd if len(cmd) < 300 else cmd[:60] + ["..."])
    print(shown)

    try:
        proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
        for line in proc.stdout:
            print(line.rstrip())
        proc.wait()
        if proc.returncode != 0:
            print(f"ffmpeg failed with return code {proc.returncode}", file=sys.stderr)
            return 3
    except FileNotFoundError:
        print("ffmpeg/ffprobe not found in PATH.", file=sys.stderr)
        return 4

    print(f"Saved: {out_path}")
    return 0

# ----------------- 入口 -----------------

def main():
    parser = argparse.ArgumentParser(description="Images→PPTX / Videos→MP4 stacking CLI.")
    sub = parser.add_subparsers(dest="cmd", required=True)

    # images 子命令保持不变……
    p_img = sub.add_parser("images", help="Linearly stack images and export PPTX (independent shapes).")
    p_img.add_argument("--files", nargs="+", required=True, help="Image files or globs.")
    p_img.add_argument("--direction", choices=["horizontal", "vertical"], default="horizontal")
    p_img.add_argument("--gap-px", type=int, default=5)
    p_img.add_argument("--border-pt", type=float, default=1.0)
    p_img.add_argument("--out", type=str, default="output.pptx")
    p_img.set_defaults(func=cmd_images)

    # videos 子命令（精简短参）
    p_vid = sub.add_parser("videos", help="Grid-stack videos into MP4 using ffmpeg xstack (strict readme).")
    p_vid.add_argument("-i", "--dir", type=str, required=True,
                       help="Folder containing videos; readme.txt assumed inside.")
    p_vid.add_argument("--rows", type=int, required=True, help="Number of rows.")
    p_vid.add_argument("--cols", type=int, required=True, help="Number of cols.")
    p_vid.add_argument("-m", "--mode", choices=["h", "v"], required=True,
                       help="h=row-major with row captions; v=col-major with col captions.")
    p_vid.add_argument("-o", "--out", type=str, required=True, help="Output MP4 path.")

    # 可选项（均有默认；用户可不写）
    p_vid.add_argument("--gap-px", type=int, default=5)
    p_vid.add_argument("--outer-border-px", type=int, default=5)
    p_vid.add_argument("--title-band-px", type=int, default=40) # 单个视频的标题与视频的间隔
    p_vid.add_argument("--rowcap-band-px", type=int, default=150) # 行标题与视频的间隔
    p_vid.add_argument("--colcap-band-px", type=int, default=20) 
    p_vid.add_argument("--title-fontsize", type=int, default=26)
    p_vid.add_argument("--rowcap-fontsize", type=int, default=30) # 行标题的大小
    p_vid.add_argument("--colcap-fontsize", type=int, default=30)
    p_vid.add_argument("--keep-audio", choices=["first", "none"], default="first")
    p_vid.add_argument("--fontfile", type=str, default="")
    p_vid.set_defaults(func=cmd_videos)

    args = parser.parse_args()
    return args.func(args)

if __name__ == "__main__":
    sys.exit(main())